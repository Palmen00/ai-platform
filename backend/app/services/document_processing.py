import csv
import json
from collections import Counter
from datetime import date
from io import BytesIO
from io import StringIO
import os
from pathlib import Path
import re
import shutil
from statistics import mean
import subprocess

from pypdf import PdfReader

from app.config import settings
from app.schemas.document import DocumentSignal
from app.services.gliner_service import GLiNEREntityService
from app.services.unstructured_service import UnstructuredPartitionService


class DocumentProcessingService:
    IMAGE_SUFFIXES = {".jpg", ".jpeg", ".png", ".bmp", ".tif", ".tiff", ".webp"}
    WORD_SUFFIXES = {".docx"}
    SPREADSHEET_SUFFIXES = {".xlsx"}
    PRESENTATION_SUFFIXES = {".pptx"}
    TEXT_SUFFIXES = {".txt", ".text", ".log", ".rst"}
    MARKDOWN_SUFFIXES = {".md", ".markdown", ".mdx"}
    JSON_SUFFIXES = {".json", ".jsonl", ".ndjson"}
    CSV_SUFFIXES = {".csv", ".tsv"}
    CONFIG_SUFFIXES = {
        ".yml", ".yaml", ".toml", ".ini", ".cfg", ".conf", ".env", ".properties", ".xml",
    }
    CODE_SUFFIXES = {
        ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".cs", ".go", ".rs", ".php", ".rb",
        ".c", ".cc", ".cpp", ".cxx", ".h", ".hpp", ".swift", ".kt", ".kts", ".scala",
        ".sh", ".bash", ".zsh", ".ps1", ".psm1", ".psd1", ".sql", ".html", ".htm",
        ".css", ".scss", ".less", ".vue", ".svelte",
    }
    TEXTLIKE_MIME_TYPES = {
        "application/json",
        "application/ld+json",
        "application/xml",
        "application/javascript",
        "application/x-javascript",
        "application/sql",
        "application/x-sh",
        "application/x-shellscript",
    }
    SIGNAL_STOPWORDS = {
        "the", "and", "for", "with", "that", "this", "from", "your", "you",
        "are", "was", "were", "have", "has", "had", "into", "their", "they",
        "there", "what", "which", "when", "where", "will", "would", "could",
        "shall", "should", "can", "may", "about", "after", "before", "between",
        "through", "under", "over", "within", "without", "onto", "into",
        "och", "att", "det", "den", "som", "har", "kan", "med", "utan",
        "från", "fran", "eller", "samt", "inga", "några", "nagra", "mina",
        "dina", "våra", "vara", "inte", "detta", "denna", "till", "hos",
        "invoice", "faktura", "salgsfaktura", "document", "documents", "page",
        "pages", "report", "roadmap", "policy", "architecture", "features",
    }
    ENTITY_NOISE_TERMS = {
        "art",
        "bekr",
        "nit",
        "org",
        "tglm",
        "node",
        "index",
        "require",
        "login",
        "sso",
        "corp",
        "diagrams",
        "linkedin",
        "github",
        "nummer",
        "datum",
        "beskrivelse",
        "description",
        "total",
        "subtotal",
        "vat",
        "ocr",
    }
    COMPANY_SUFFIXES = {
        "ab": "AB",
        "aps": "ApS",
        "as": "AS",
        "bv": "BV",
        "gmbh": "GmbH",
        "inc": "Inc",
        "llc": "LLC",
        "ltd": "Ltd",
        "oy": "OY",
        "sa": "SA",
        "sarl": "SARL",
    }
    DOCUMENT_TYPE_RULES = {
        "invoice": (
            ("invoice", 3.0),
            ("faktura", 2.8),
            ("salgsfaktura", 3.0),
            ("invoice no", 1.8),
            ("invoice number", 1.8),
            ("invoice date", 1.8),
            ("fakturadato", 1.8),
            ("bill to", 1.2),
            ("amount due", 1.5),
            ("payment terms", 1.0),
            ("subtotal", 0.8),
            ("vat", 0.7),
            ("due date", 0.6),
        ),
        "contract": (
            ("agreement", 2.2),
            ("contract", 2.6),
            ("effective date", 1.9),
            ("party", 0.8),
            ("parties", 1.0),
            ("terms and conditions", 1.1),
            ("signed", 1.0),
            ("signature", 0.8),
        ),
        "insurance": (
            ("insurance", 2.8),
            ("policy", 2.0),
            ("insured", 1.3),
            ("coverage", 1.3),
            ("beneficiary", 1.2),
            ("claim", 1.0),
            ("premium", 0.8),
        ),
        "policy": (
            ("policy", 2.1),
            ("purpose", 0.8),
            ("scope", 0.9),
            ("compliance", 1.1),
            ("applies to", 0.8),
            ("responsibilities", 0.8),
        ),
        "roadmap": (
            ("roadmap", 3.0),
            ("week 1", 1.2),
            ("milestone", 1.2),
            ("timeline", 1.0),
            ("phase", 0.8),
        ),
        "architecture": (
            ("architecture", 3.0),
            ("system layers", 1.8),
            ("backend responsibilities", 1.8),
            ("api layer", 1.3),
        ),
        "report": (
            ("report", 2.2),
            ("summary", 0.8),
            ("findings", 1.0),
            ("analysis", 1.0),
            ("recommendation", 0.8),
        ),
        "form": (
            ("form", 1.2),
            ("first name", 1.4),
            ("last name", 1.4),
            ("date of birth", 1.4),
            ("address", 0.7),
            ("signature", 0.7),
        ),
        "receipt": (
            ("receipt", 3.0),
            ("transaction", 0.9),
            ("payment method", 1.2),
            ("thank you for your purchase", 1.4),
            ("cashier", 1.0),
        ),
        "quote": (
            ("quotation", 2.8),
            ("quote", 2.6),
            ("estimate", 1.8),
            ("valid until", 0.9),
        ),
        "features": (
            ("current features", 2.6),
            ("implemented", 0.8),
            ("model status", 0.8),
            ("available models", 0.8),
        ),
        "code": (
            ("function", 0.8),
            ("class", 0.8),
            ("import", 0.5),
            ("export", 0.5),
            ("return", 0.4),
            ("def ", 0.8),
            ("const ", 0.5),
            ("async ", 0.4),
        ),
        "config": (
            ("version", 0.5),
            ("services", 0.7),
            ("environment", 0.6),
            ("apiversion", 0.6),
            ("database", 0.5),
            ("host", 0.3),
            ("port", 0.3),
        ),
    }
    DATE_KIND_LABELS = {
        "invoice_date": "Invoice date",
        "issue_date": "Issue date",
        "effective_date": "Effective date",
        "signed_date": "Signed date",
        "policy_date": "Policy date",
        "due_date": "Due date",
        "document_date": "Document date",
        "filename_date": "Filename date",
    }
    MONTH_LOOKUP = {
        "jan": 1,
        "january": 1,
        "januari": 1,
        "feb": 2,
        "february": 2,
        "februari": 2,
        "mar": 3,
        "march": 3,
        "mars": 3,
        "apr": 4,
        "april": 4,
        "may": 5,
        "maj": 5,
        "jun": 6,
        "june": 6,
        "juni": 6,
        "jul": 7,
        "july": 7,
        "juli": 7,
        "aug": 8,
        "august": 8,
        "sep": 9,
        "sept": 9,
        "september": 9,
        "oct": 10,
        "october": 10,
        "okt": 10,
        "oktober": 10,
        "nov": 11,
        "november": 11,
        "dec": 12,
        "december": 12,
    }

    def __init__(self) -> None:
        self._docker_available: bool | None = None
        self._ocrmypdf_image_ready: bool | None = None

    def extract_text(self, file_path: Path, content_type: str) -> str:
        result = self.extract_document(file_path, content_type)
        return str(result.get("text", ""))

    def extract_document(
        self,
        file_path: Path,
        content_type: str,
    ) -> dict[str, object]:
        suffix = file_path.suffix.lower()

        if suffix in self.WORD_SUFFIXES:
            return {
                "text": self._extract_docx_file(file_path),
                "ocr_used": False,
                "ocr_status": "not_needed",
                "ocr_error": None,
            }

        if suffix in self.SPREADSHEET_SUFFIXES:
            return {
                "text": self._extract_xlsx_file(file_path),
                "ocr_used": False,
                "ocr_status": "not_needed",
                "ocr_error": None,
            }

        if suffix in self.PRESENTATION_SUFFIXES:
            return {
                "text": self._extract_pptx_file(file_path),
                "ocr_used": False,
                "ocr_status": "not_needed",
                "ocr_error": None,
            }

        if self._is_text_like_document(suffix=suffix, content_type=content_type):
            return {
                "text": self._extract_text_file(file_path, suffix),
                "ocr_used": False,
                "ocr_status": "not_needed",
                "ocr_error": None,
            }

        if suffix == ".pdf" or content_type == "application/pdf":
            return self._extract_pdf(file_path)

        if suffix in self.IMAGE_SUFFIXES or content_type.startswith("image/"):
            return self._extract_image(file_path)

        raise ValueError(f"Unsupported document type: {suffix or content_type}")

    def chunk_text(
        self,
        text: str,
        document_name: str = "",
        content_type: str = "",
        file_path: Path | None = None,
    ) -> list[dict[str, str | int | None]]:
        sections = self.segment_sections(
            text,
            document_name=document_name,
            content_type=content_type,
            file_path=file_path,
        )
        if not sections:
            return []

        chunk_size = settings.document_chunk_size
        overlap = min(settings.document_chunk_overlap, chunk_size // 2)
        chunks: list[dict[str, str | int | None]] = []
        global_offset = 0

        for section in sections:
            section_title = str(section.get("title") or "").strip() or None
            page_number = section.get("page_number")
            normalized_content = " ".join(str(section.get("content", "")).split())
            if not normalized_content:
                continue

            start = 0
            while start < len(normalized_content):
                end = min(start + chunk_size, len(normalized_content))
                body = normalized_content[start:end].strip()
                if not body:
                    break

                content = body
                if section_title and not body.lower().startswith(section_title.lower()):
                    content = f"{section_title}\n\n{body}"

                chunks.append(
                    {
                        "index": len(chunks),
                        "start": global_offset + start,
                        "end": global_offset + end,
                        "content": content,
                        "section_title": section_title,
                        "page_number": int(page_number) if page_number else None,
                        "source_kind": self.detect_source_kind(
                            document_name=document_name,
                            content_type=content_type,
                        ),
                    }
                )

                if end >= len(normalized_content):
                    break

                start = max(end - overlap, start + 1)

            global_offset += len(normalized_content) + 2

        return chunks

    def detect_document_title(
        self,
        text: str,
        document_name: str,
        content_type: str = "",
        file_path: Path | None = None,
    ) -> str:
        if file_path is not None:
            title = UnstructuredPartitionService().extract_document_title(
                file_path=file_path,
                content_type=content_type,
            )
            if title:
                return title

        lines = [line.strip() for line in text.splitlines() if line.strip()]
        for line in lines[:8]:
            if self._is_section_heading(line):
                return self._clean_heading(line)

        return Path(document_name).stem.replace("_", " ").replace("-", " ").strip()

    def detect_source_kind(self, document_name: str, content_type: str) -> str:
        suffix = Path(document_name).suffix.lower()
        if suffix == ".pdf" or content_type == "application/pdf":
            return "pdf"
        if suffix in self.WORD_SUFFIXES:
            return "word"
        if suffix in self.SPREADSHEET_SUFFIXES:
            return "spreadsheet"
        if suffix in self.PRESENTATION_SUFFIXES:
            return "presentation"
        if suffix in self.MARKDOWN_SUFFIXES:
            return "markdown"
        if suffix in self.JSON_SUFFIXES:
            return "json"
        if suffix in self.CSV_SUFFIXES:
            return "csv"
        if suffix in self.CONFIG_SUFFIXES:
            return "config"
        if suffix in self.IMAGE_SUFFIXES or content_type.startswith("image/"):
            return "image"
        if suffix in self.CODE_SUFFIXES:
            return "code"
        if suffix in self.TEXT_SUFFIXES or self._is_text_like_content_type(content_type):
            return "text"
        return "document"

    def detect_document_type(
        self,
        text: str,
        document_name: str,
        content_type: str,
    ) -> str:
        normalized_name = Path(document_name).stem.lower().replace("_", " ").replace("-", " ")
        normalized_text = " ".join(self._normalize_extracted_text(text).lower().split())
        haystack = f"{normalized_name}\n{normalized_text[:16000]}"
        source_kind = self.detect_source_kind(document_name, content_type)
        best_type = "document"
        best_score = 0.0

        for document_type, rules in self.DOCUMENT_TYPE_RULES.items():
            score = 0.0
            for phrase, weight in rules:
                if phrase in haystack:
                    score += weight

            if document_type in normalized_name:
                score += 1.2

            if score > best_score:
                best_score = score
                best_type = document_type

        if source_kind == "presentation":
            return "presentation"

        if source_kind == "spreadsheet" and best_score < 3.5:
            return "spreadsheet"

        if source_kind == "word" and best_type == "document":
            return "word"

        if best_score < 1.5:
            if source_kind == "image" and normalized_text:
                return "form"
            if source_kind == "code":
                return "code"
            if source_kind == "config":
                return "config"
            return "document"

        return best_type

    def detect_document_date(
        self,
        text: str,
        document_name: str,
        document_type: str | None = None,
    ) -> tuple[str | None, str | None, str | None]:
        candidates = self._extract_date_candidates(text=text, document_name=document_name)
        if not candidates:
            return (None, None, None)

        ranked_candidates = sorted(
            candidates,
            key=lambda item: (
                self._score_date_candidate(
                    candidate_kind=str(item["kind"]),
                    document_type=document_type or "document",
                ),
                -int(item["position"]),
            ),
            reverse=True,
        )
        best = ranked_candidates[0]
        iso_value = str(best["date"])
        date_kind = str(best["kind"])
        date_label = f"{self.DATE_KIND_LABELS.get(date_kind, 'Document date')}: {best['raw']}"
        return (iso_value, date_label, date_kind)

    def detect_document_entities(
        self,
        text: str,
        document_name: str,
        document_type: str | None = None,
    ) -> list[str]:
        normalized_text = self._normalize_extracted_text(text)
        lines = [line.strip() for line in normalized_text.splitlines() if line.strip()]
        candidates: list[str] = []

        labeled_patterns = (
            r"^(?:vendor|supplier|seller|company|customer|bill to|ship to|invoice to|leverant[oö]r|kund)\s*[:\-]\s*(.+)$",
        )
        for line in lines[:80]:
            for pattern in labeled_patterns:
                match = re.match(pattern, line, flags=re.IGNORECASE)
                if not match:
                    continue
                normalized = self._normalize_entity_name(match.group(1))
                if normalized:
                    candidates.append(normalized)

            line_candidate = self._line_entity_candidate(line, document_type=document_type)
            if line_candidate:
                candidates.append(line_candidate)

        domain_entities = re.findall(
            r"\b(?:https?://)?(?:www\.)?([a-z0-9][a-z0-9-]{1,}\.[a-z]{2,})\b",
            normalized_text.lower(),
        )
        for domain in domain_entities[:12]:
            domain_candidate = self._entity_from_domain(domain)
            if domain_candidate:
                candidates.append(domain_candidate)

        filename_candidate = self._normalize_entity_name(Path(document_name).stem)
        if filename_candidate and any(token.isalpha() for token in filename_candidate.split()):
            if document_type in {"invoice", "contract", "insurance", "quote"}:
                candidates.append(filename_candidate)

        gliner_candidates = GLiNEREntityService().extract_candidate_entities(
            normalized_text,
            document_type=document_type,
        )
        for candidate_text, label, _score in gliner_candidates:
            normalized_candidate = self._normalize_entity_name(candidate_text)
            if normalized_candidate and not self._is_noise_entity(normalized_candidate):
                candidates.append(normalized_candidate)

        deduplicated: list[str] = []
        seen: set[str] = set()
        for candidate in candidates:
            canonical = self._entity_match_key(candidate)
            if not canonical or canonical in seen:
                continue
            seen.add(canonical)
            deduplicated.append(candidate)

        return deduplicated[:8]

    def detect_document_signals(
        self,
        text: str,
        document_name: str,
        document_type: str | None = None,
        document_title: str | None = None,
        document_entities: list[str] | None = None,
    ) -> list[DocumentSignal]:
        normalized_text = self._normalize_extracted_text(text)
        signals: dict[tuple[str, str], DocumentSignal] = {}

        def upsert_signal(
            value: str,
            *,
            category: str,
            score: float,
            source: str,
        ) -> None:
            normalized = self._signal_key(value)
            if not normalized:
                return
            if self._is_noise_signal_value(value, category=category):
                return
            key = (category, normalized)
            existing = signals.get(key)
            if existing is None or float(existing.score) < score:
                signals[key] = DocumentSignal(
                    value=value.strip(),
                    normalized=normalized,
                    category=category,
                    score=round(score, 4),
                    source=source,
                )

        for entity in document_entities or []:
            upsert_signal(entity, category="entity", score=0.96, source="entity")

        if document_title:
            upsert_signal(document_title, category="title", score=0.88, source="title")

        filename_title = Path(document_name).stem.replace("_", " ").replace("-", " ").strip()
        if filename_title:
            upsert_signal(filename_title, category="filename", score=0.72, source="filename")

        sections = self.segment_sections(text, document_name=document_name)
        for section in sections[:12]:
            section_title = str(section.get("title") or "").strip()
            if section_title:
                upsert_signal(section_title, category="section", score=0.76, source="section")

        text_window = normalized_text[:40000]
        token_counts = self._signal_token_counts(text_window)
        max_count = max(token_counts.values(), default=1)
        for token, count in token_counts.most_common(24):
            base_score = 0.24 + min(count / max_count, 1.0) * 0.42
            if document_type and token == document_type:
                base_score += 0.16
            upsert_signal(token, category="keyword", score=min(base_score, 0.84), source="text")

        phrase_counts = self._signal_phrase_counts(text_window)
        max_phrase_count = max(phrase_counts.values(), default=1)
        for phrase, count in phrase_counts.most_common(18):
            base_score = 0.3 + min(count / max_phrase_count, 1.0) * 0.36
            upsert_signal(phrase, category="phrase", score=min(base_score, 0.82), source="text")

        ranked_signals = sorted(
            signals.values(),
            key=lambda item: (float(item.score), len(str(item.value))),
            reverse=True,
        )
        return ranked_signals[:32]

    def segment_sections(
        self,
        text: str,
        document_name: str = "",
        content_type: str = "",
        file_path: Path | None = None,
    ) -> list[dict[str, str | int | None]]:
        lines = [self._normalize_line(line) for line in text.splitlines()]
        fallback_title = Path(document_name).stem.replace("_", " ").replace("-", " ").strip()
        if file_path is not None:
            unstructured_sections = UnstructuredPartitionService().extract_sections(
                file_path=file_path,
                content_type=content_type,
                fallback_title=fallback_title,
            )
            if unstructured_sections:
                return unstructured_sections

        fallback_title = self.detect_document_title(text, document_name)
        sections: list[dict[str, str | int | None]] = []
        current_title = fallback_title or "Document"
        current_page: int | None = None
        buffer: list[str] = []

        def flush_section() -> None:
            normalized_content = "\n".join(line for line in buffer if line.strip()).strip()
            if not normalized_content:
                return

            sections.append(
                {
                    "title": current_title,
                    "content": normalized_content,
                    "page_number": current_page,
                }
            )

        for raw_line in lines:
            line = raw_line.strip()
            if not line:
                if buffer and buffer[-1] != "":
                    buffer.append("")
                continue

            page_match = re.fullmatch(r"\[page\s+(\d+)\]", line.lower())
            if page_match:
                flush_section()
                buffer = []
                current_page = int(page_match.group(1))
                continue

            if self._is_ordered_list_item(line):
                buffer.append(line)
                continue

            if self._is_section_heading(line):
                flush_section()
                buffer = []
                current_title = self._clean_heading(line)
                continue

            buffer.append(line)

        flush_section()

        if sections:
            return sections

        normalized_text = "\n".join(line for line in lines if line.strip()).strip()
        if not normalized_text:
            return []

        return [{"title": fallback_title or "Document", "content": normalized_text, "page_number": None}]

    def _is_text_like_document(self, suffix: str, content_type: str) -> bool:
        if suffix in (
            self.TEXT_SUFFIXES
            | self.MARKDOWN_SUFFIXES
            | self.JSON_SUFFIXES
            | self.CSV_SUFFIXES
            | self.CONFIG_SUFFIXES
            | self.CODE_SUFFIXES
        ):
            return True

        return self._is_text_like_content_type(content_type)

    def _is_text_like_content_type(self, content_type: str) -> bool:
        normalized = (content_type or "").split(";", 1)[0].strip().lower()
        if not normalized:
            return False
        if normalized.startswith("text/"):
            return True
        if normalized in self.TEXTLIKE_MIME_TYPES:
            return True
        if normalized.endswith("+json") or normalized.endswith("+xml"):
            return True
        return False

    def _extract_text_file(self, file_path: Path, suffix: str) -> str:
        raw_text = file_path.read_text(encoding="utf-8", errors="ignore")

        if suffix in self.JSON_SUFFIXES:
            try:
                payload = json.loads(raw_text)
                return json.dumps(payload, ensure_ascii=False, indent=2)
            except json.JSONDecodeError:
                return raw_text

        if suffix in self.CSV_SUFFIXES:
            reader = csv.reader(StringIO(raw_text))
            rows = [" | ".join(cell.strip() for cell in row) for row in reader]
            return "\n".join(rows)

        return raw_text

    def _extract_docx_file(self, file_path: Path) -> str:
        from docx import Document

        document = Document(file_path)
        parts: list[str] = []

        for paragraph in document.paragraphs:
            text = " ".join(paragraph.text.split()).strip()
            if not text:
                continue
            style_name = getattr(getattr(paragraph, "style", None), "name", "") or ""
            if style_name.lower().startswith("heading"):
                parts.append(text)
                parts.append("")
            else:
                parts.append(text)

        for table in document.tables:
            parts.append("")
            parts.append("Table")
            for row in table.rows:
                values = [
                    " ".join(cell.text.split()).strip()
                    for cell in row.cells
                    if " ".join(cell.text.split()).strip()
                ]
                if values:
                    parts.append(" | ".join(values))

        return "\n".join(parts).strip()

    def _extract_xlsx_file(self, file_path: Path) -> str:
        from openpyxl import load_workbook

        workbook = load_workbook(file_path, data_only=True)
        parts: list[str] = []

        for sheet in workbook.worksheets:
            parts.append(f"# Sheet {sheet.title}")
            rows = list(sheet.iter_rows(values_only=True))
            headers: list[str] = []
            if rows:
                headers = [
                    " ".join(str(cell).split()).strip()
                    for cell in rows[0]
                    if cell is not None and " ".join(str(cell).split()).strip()
                ]

            for row_index, row in enumerate(rows[1:] if headers else rows, start=1):
                values = [
                    " ".join(str(cell).split()).strip()
                    if cell is not None
                    else ""
                    for cell in row
                ]
                if not any(values):
                    continue

                row_pairs: list[str] = []
                for column_index, value in enumerate(values):
                    if not value:
                        continue
                    label = (
                        headers[column_index]
                        if column_index < len(headers) and headers[column_index]
                        else f"Column {column_index + 1}"
                    )
                    row_pairs.append(f"{label}: {value}")

                if not row_pairs:
                    continue

                section_title = (
                    values[0]
                    if values and values[0]
                    else f"Row {row_index}"
                )
                parts.append(f"## {sheet.title} / {section_title}")
                parts.append(" ; ".join(row_pairs))
            parts.append("")

        return "\n".join(parts).strip()

    def _extract_pptx_file(self, file_path: Path) -> str:
        from pptx import Presentation

        presentation = Presentation(file_path)
        parts: list[str] = []

        for index, slide in enumerate(presentation.slides, start=1):
            title_text = ""
            if slide.shapes.title and slide.shapes.title.text:
                title_text = " ".join(slide.shapes.title.text.split()).strip()

            heading = title_text or f"Slide {index}"
            parts.append(f"# Slide {index}: {heading}")

            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = " ".join(str(shape.text).split()).strip()
                if not text or text == title_text:
                    continue
                parts.append(f"- {text}")

            parts.append("")

        return "\n".join(parts).strip()

    def _extract_pdf(self, file_path: Path) -> dict[str, object]:
        direct_text = self._extract_pdf_text(file_path)
        if self._has_meaningful_text(direct_text):
            return {
                "text": direct_text,
                "ocr_used": False,
                "ocr_status": "not_needed",
                "ocr_engine": None,
                "ocr_error": None,
            }

        if not settings.ocr_enabled:
            return {
                "text": direct_text,
                "ocr_used": False,
                "ocr_status": "unavailable",
                "ocr_engine": None,
                "ocr_error": "OCR fallback is disabled in the current runtime configuration.",
            }

        ocrmypdf_error: str | None = None
        if settings.ocrmypdf_enabled:
            try:
                ocrmypdf_text = self._extract_pdf_with_ocrmypdf(file_path)
            except RuntimeError as exc:
                ocrmypdf_error = str(exc)
            except Exception as exc:
                ocrmypdf_error = f"OCRmyPDF failed: {exc}"
            else:
                if self._has_meaningful_text(ocrmypdf_text):
                    return {
                        "text": self._cleanup_ocr_text(ocrmypdf_text),
                        "ocr_used": True,
                        "ocr_status": "used",
                        "ocr_engine": "ocrmypdf",
                        "ocr_error": None,
                    }

        try:
            ocr_text = self._extract_pdf_with_ocr(file_path)
        except RuntimeError as exc:
            return {
                "text": direct_text,
                "ocr_used": False,
                "ocr_status": "unavailable",
                "ocr_engine": None,
                "ocr_error": self._merge_ocr_errors(ocrmypdf_error, str(exc)),
            }
        except Exception as exc:
            return {
                "text": direct_text,
                "ocr_used": False,
                "ocr_status": "failed",
                "ocr_engine": None,
                "ocr_error": self._merge_ocr_errors(ocrmypdf_error, f"OCR failed: {exc}"),
            }

        if self._has_meaningful_text(ocr_text):
            return {
                "text": self._cleanup_ocr_text(ocr_text),
                "ocr_used": True,
                "ocr_status": "used",
                "ocr_engine": "tesseract",
                "ocr_error": None,
            }

        return {
            "text": direct_text,
            "ocr_used": False,
            "ocr_status": "failed",
            "ocr_engine": None,
            "ocr_error": self._merge_ocr_errors(
                ocrmypdf_error,
                "OCR ran but did not find enough readable text.",
            ),
        }

    def _extract_image(self, file_path: Path) -> dict[str, object]:
        if not settings.ocr_enabled:
            return {
                "text": "",
                "ocr_used": False,
                "ocr_status": "unavailable",
                "ocr_engine": None,
                "ocr_error": "OCR fallback is disabled in the current runtime configuration.",
            }

        try:
            ocr_text = self._extract_image_file_with_ocr(file_path)
        except RuntimeError as exc:
            return {
                "text": "",
                "ocr_used": False,
                "ocr_status": "unavailable",
                "ocr_engine": None,
                "ocr_error": str(exc),
            }
        except Exception as exc:
            return {
                "text": "",
                "ocr_used": False,
                "ocr_status": "failed",
                "ocr_engine": None,
                "ocr_error": f"OCR failed: {exc}",
            }

        if self._has_meaningful_text(ocr_text):
            return {
                "text": self._cleanup_ocr_text(ocr_text),
                "ocr_used": True,
                "ocr_status": "used",
                "ocr_engine": "tesseract",
                "ocr_error": None,
            }

        return {
            "text": "",
            "ocr_used": False,
            "ocr_status": "failed",
            "ocr_engine": None,
            "ocr_error": "OCR ran but did not find enough readable text.",
        }

    def _extract_pdf_text(self, file_path: Path) -> str:
        reader = PdfReader(str(file_path))
        pages = []
        for index, page in enumerate(reader.pages, start=1):
            text = self._normalize_extracted_text(page.extract_text() or "")
            if not text.strip():
                continue

            pages.append(f"[Page {index}]\n{text.strip()}")

        return "\n\n".join(pages)

    def _extract_pdf_with_ocrmypdf(self, file_path: Path) -> str:
        if not settings.ocrmypdf_use_docker:
            raise RuntimeError("OCRmyPDF is enabled but Docker execution is disabled.")

        self._ensure_ocrmypdf_runtime()

        output_dir = settings.ocrmypdf_cache_dir
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"{file_path.stem}.ocrmypdf.pdf"
        if output_path.exists():
            output_path.unlink()

        input_dir = file_path.parent.resolve()
        output_dir_resolved = output_dir.resolve()
        language = self._resolve_ocrmypdf_language()
        command = [
            "docker",
            "run",
            "--rm",
            "-v",
            f"{input_dir}:/input:ro",
            "-v",
            f"{output_dir_resolved}:/output",
            settings.ocrmypdf_docker_image,
            "--force-ocr",
            "--skip-big",
            str(settings.ocrmypdf_skip_big),
            "--optimize",
            "0",
            "--output-type",
            "pdf",
            "-l",
            language,
            f"/input/{file_path.name}",
            f"/output/{output_path.name}",
        ]

        try:
            subprocess.run(
                command,
                check=True,
                capture_output=True,
                text=True,
                timeout=settings.ocrmypdf_timeout_seconds,
            )
        except FileNotFoundError as exc:
            self._docker_available = False
            raise RuntimeError("OCRmyPDF Docker fallback requires Docker to be installed.") from exc
        except subprocess.TimeoutExpired as exc:
            raise RuntimeError("OCRmyPDF timed out while preprocessing the PDF.") from exc
        except subprocess.CalledProcessError as exc:
            stderr = (exc.stderr or "").strip()
            stdout = (exc.stdout or "").strip()
            detail = stderr or stdout or "unknown OCRmyPDF error"
            raise RuntimeError(f"OCRmyPDF preprocessing failed: {detail}") from exc

        if not output_path.exists():
            raise RuntimeError("OCRmyPDF did not produce an output PDF.")

        try:
            return self._extract_pdf_text(output_path)
        finally:
            try:
                output_path.unlink()
            except FileNotFoundError:
                pass

    def _extract_pdf_with_ocr(self, file_path: Path) -> str:
        try:
            import fitz
        except ImportError as exc:
            raise RuntimeError(
                "OCR support requires PyMuPDF to render scanned PDF pages."
            ) from exc

        try:
            import pytesseract
            from PIL import Image, ImageFilter, ImageOps
        except ImportError as exc:
            raise RuntimeError(
                "OCR support requires pytesseract and Pillow."
            ) from exc

        ocr_language, tessdata_dir = self._prepare_ocr_runtime(pytesseract)

        document = fitz.open(file_path)
        pages: list[str] = []
        matrix = fitz.Matrix(settings.ocr_render_scale, settings.ocr_render_scale)

        for index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image = Image.open(BytesIO(pixmap.tobytes("png")))
            text = self._extract_best_ocr_text_from_image(
                image=image,
                ocr_language=ocr_language,
                tessdata_dir=tessdata_dir,
                pytesseract_module=pytesseract,
                image_filter_module=ImageFilter,
                image_ops_module=ImageOps,
            )
            normalized_text = self._normalize_extracted_text(text)
            if not normalized_text.strip():
                continue

            pages.append(f"[Page {index}]\n{normalized_text.strip()}")

        return "\n\n".join(pages)

    def _ensure_ocrmypdf_runtime(self) -> None:
        if not self._docker_is_available():
            raise RuntimeError("OCRmyPDF Docker fallback requires Docker to be running.")

        if self._ocrmypdf_image_ready is True:
            return

        image_name = settings.ocrmypdf_docker_image
        inspect_command = ["docker", "image", "inspect", image_name]
        inspected = subprocess.run(
            inspect_command,
            capture_output=True,
            text=True,
            check=False,
        )
        if inspected.returncode == 0:
            self._ocrmypdf_image_ready = True
            return

        if not settings.ocrmypdf_auto_build:
            raise RuntimeError(
                f"OCRmyPDF Docker image '{image_name}' is not available locally."
            )

        dockerfile_path = settings.repo_root / "infra" / "ocrmypdf" / "Dockerfile"
        if not dockerfile_path.exists():
            raise RuntimeError(
                "OCRmyPDF helper image is missing its Dockerfile."
            )

        try:
            subprocess.run(
                [
                    "docker",
                    "build",
                    "-t",
                    image_name,
                    "-f",
                    str(dockerfile_path),
                    str(settings.repo_root),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=max(settings.ocrmypdf_timeout_seconds, 600),
            )
        except subprocess.CalledProcessError as exc:
            stderr = (exc.stderr or "").strip()
            stdout = (exc.stdout or "").strip()
            detail = stderr or stdout or "unknown docker build error"
            raise RuntimeError(f"Could not build the OCRmyPDF helper image: {detail}") from exc
        except subprocess.TimeoutExpired as exc:
            raise RuntimeError("Building the OCRmyPDF helper image timed out.") from exc

        self._ocrmypdf_image_ready = True

    def _docker_is_available(self) -> bool:
        if self._docker_available is not None:
            return self._docker_available

        try:
            result = subprocess.run(
                ["docker", "version", "--format", "{{.Server.Version}}"],
                check=False,
                capture_output=True,
                text=True,
                timeout=20,
            )
        except (FileNotFoundError, subprocess.TimeoutExpired):
            self._docker_available = False
            return False

        self._docker_available = result.returncode == 0
        return self._docker_available

    def _resolve_ocrmypdf_language(self) -> str:
        configured = [item.strip() for item in settings.ocr_language.split("+") if item.strip()]
        mapped: list[str] = []
        for language in configured:
            lowered = language.lower()
            if lowered in {"eng", "en"} and "eng" not in mapped:
                mapped.append("eng")
            elif lowered in {"swe", "sv"} and "swe" not in mapped:
                mapped.append("swe")
            elif lowered and lowered not in mapped:
                mapped.append(lowered)
        if not mapped:
            mapped.append("eng")
        return "+".join(mapped)

    def _merge_ocr_errors(self, *errors: str | None) -> str:
        parts: list[str] = []
        seen: set[str] = set()
        for error in errors:
            normalized = str(error or "").strip()
            if not normalized or normalized in seen:
                continue
            seen.add(normalized)
            parts.append(normalized)
        return " ".join(parts)

    def _extract_image_file_with_ocr(self, file_path: Path) -> str:
        try:
            import pytesseract
            from PIL import Image, ImageFilter, ImageOps
        except ImportError as exc:
            raise RuntimeError(
                "OCR support requires pytesseract and Pillow."
            ) from exc

        ocr_language, tessdata_dir = self._prepare_ocr_runtime(pytesseract)
        image = Image.open(file_path)
        oriented = self._auto_orient_image(image, pytesseract)

        return self._extract_best_ocr_text_from_image(
            image=oriented,
            ocr_language=ocr_language,
            tessdata_dir=tessdata_dir,
            pytesseract_module=pytesseract,
            image_filter_module=ImageFilter,
            image_ops_module=ImageOps,
        )

    def _has_meaningful_text(self, text: str) -> bool:
        normalized = self._normalize_extracted_text(text)
        alphanumeric_count = len(re.findall(r"[A-Za-z0-9]", normalized))
        words = re.findall(r"[A-Za-z][A-Za-z'-]{2,}", normalized)
        unique_words = len(set(word.lower() for word in words))
        return (
            alphanumeric_count >= settings.ocr_min_characters
            or (alphanumeric_count >= 24 and unique_words >= 3)
        )

    def _extract_best_ocr_text_from_image(
        self,
        image: object,
        ocr_language: str,
        tessdata_dir: str | None,
        pytesseract_module: object,
        image_filter_module: object,
        image_ops_module: object,
    ) -> str:
        variants = self._build_ocr_variants(
            image=image,
            image_filter_module=image_filter_module,
            image_ops_module=image_ops_module,
        )
        candidates: list[tuple[float, float, str]] = []

        for variant in variants:
            for psm in settings.ocr_psm_modes:
                score, confidence, normalized_text = self._extract_ocr_candidate(
                    image=variant,
                    psm=psm,
                    ocr_language=ocr_language,
                    tessdata_dir=tessdata_dir,
                    pytesseract_module=pytesseract_module,
                )
                if score <= 0:
                    continue

                candidates.append((score, confidence, normalized_text))

        if not candidates:
            return ""

        candidates.sort(key=lambda item: (item[0], item[1]), reverse=True)
        return candidates[0][2]

    def _extract_ocr_candidate(
        self,
        image: object,
        psm: int,
        ocr_language: str,
        tessdata_dir: str | None,
        pytesseract_module: object,
    ) -> tuple[float, float, str]:
        try:
            raw_text = pytesseract_module.image_to_string(
                image,
                lang=ocr_language,
                config=self._build_tesseract_config(psm=psm, tessdata_dir=tessdata_dir),
            )
        except Exception:
            return (0.0, 0.0, "")

        normalized_text = self._normalize_extracted_text(raw_text)
        score = self._score_ocr_candidate(normalized_text)
        if score <= 0:
            return (0.0, 0.0, normalized_text)

        confidence = self._extract_ocr_confidence(
            image=image,
            psm=psm,
            ocr_language=ocr_language,
            tessdata_dir=tessdata_dir,
            pytesseract_module=pytesseract_module,
        )
        boosted_score = score + (confidence * 0.02)
        return (boosted_score, confidence, normalized_text)

    def _build_ocr_variants(
        self,
        image: object,
        image_filter_module: object,
        image_ops_module: object,
    ) -> list[object]:
        grayscale = image.convert("L")
        autocontrast = image_ops_module.autocontrast(grayscale)
        sharpened = autocontrast.filter(image_filter_module.SHARPEN)
        median = autocontrast.filter(image_filter_module.MedianFilter(size=3))
        thresholded = autocontrast.point(
            lambda pixel: 255 if pixel > self._image_threshold(autocontrast) else 0
        )
        upscaled = autocontrast.resize(
            (
                max(int(getattr(autocontrast, "width", 1) * 1.5), 1),
                max(int(getattr(autocontrast, "height", 1) * 1.5), 1),
            )
        )

        variants = [image, grayscale, autocontrast, sharpened, median, thresholded, upscaled]
        deduplicated: list[object] = []
        seen_sizes: set[tuple[int, int, str]] = set()

        for variant in variants:
            key = (
                int(getattr(variant, "width", 0)),
                int(getattr(variant, "height", 0)),
                str(getattr(variant, "mode", "")),
            )
            if key in seen_sizes and variant in deduplicated:
                continue
            seen_sizes.add(key)
            deduplicated.append(variant)

        return deduplicated

    def _image_threshold(self, image: object) -> int:
        try:
            histogram = image.histogram()
            pixel_count = sum(histogram)
            if pixel_count <= 0:
                return 180

            weighted_sum = sum(index * count for index, count in enumerate(histogram))
            average = weighted_sum / pixel_count
            return int(max(140, min(average, 210)))
        except Exception:
            return 180

    def _extract_ocr_confidence(
        self,
        image: object,
        psm: int,
        ocr_language: str,
        tessdata_dir: str | None,
        pytesseract_module: object,
    ) -> float:
        try:
            data = pytesseract_module.image_to_data(
                image,
                lang=ocr_language,
                config=self._build_tesseract_config(psm=psm, tessdata_dir=tessdata_dir),
                output_type=pytesseract_module.Output.DICT,
            )
        except Exception:
            return 0.0

        confidences: list[float] = []
        for confidence, text in zip(data.get("conf", []), data.get("text", [])):
            value = str(text).strip()
            if not value:
                continue

            try:
                numeric_confidence = float(confidence)
            except Exception:
                continue

            if numeric_confidence >= 0:
                confidences.append(numeric_confidence)

        if not confidences:
            return 0.0

        return mean(confidences)

    def _prepare_ocr_runtime(self, pytesseract_module: object) -> tuple[str, str | None]:
        tesseract_cmd = self._resolve_tesseract_cmd()
        if tesseract_cmd:
            pytesseract_module.pytesseract.tesseract_cmd = tesseract_cmd

        try:
            pytesseract_module.get_tesseract_version()
        except Exception as exc:
            raise RuntimeError(
                "OCR support requires the Tesseract OCR engine to be installed and reachable."
            ) from exc

        tessdata_dir = self._resolve_tessdata_dir()
        if tessdata_dir:
            os.environ["TESSDATA_PREFIX"] = tessdata_dir
        ocr_language = self._resolve_ocr_language(
            pytesseract_module=pytesseract_module,
            tessdata_dir=tessdata_dir,
        )
        return (ocr_language, tessdata_dir)

    def _build_tesseract_config(self, psm: int, tessdata_dir: str | None) -> str:
        return f"--psm {psm}"

    def _score_ocr_candidate(self, text: str) -> float:
        normalized = " ".join(text.split())
        if not normalized:
            return 0.0

        alphanumeric_count = len(re.findall(r"[A-Za-z0-9]", normalized))
        words = re.findall(r"[A-Za-z][A-Za-z'-]{2,}", normalized)
        unique_words = len(set(word.lower() for word in words))
        if alphanumeric_count < 16 and unique_words < 3:
            return 0.0
        average_word_length = mean([len(word) for word in words]) if words else 0.0
        punctuation_noise = len(re.findall(r"[^A-Za-z0-9\s.,;:!?()/%€$'-]", normalized))

        return (
            (alphanumeric_count * 0.01)
            + min(unique_words, 120) * 0.03
            + min(average_word_length, 8) * 0.2
            - min(punctuation_noise, 50) * 0.05
        )

    def _auto_orient_image(self, image: object, pytesseract_module: object) -> object:
        try:
            osd = pytesseract_module.image_to_osd(image, lang="osd")
        except Exception:
            return image

        rotation_match = re.search(r"Rotate:\s+(\d+)", str(osd))
        if not rotation_match:
            return image

        rotation = int(rotation_match.group(1)) % 360
        if rotation == 0:
            return image

        try:
            return image.rotate(-rotation, expand=True)
        except Exception:
            return image

    def _resolve_ocr_language(
        self,
        pytesseract_module: object,
        tessdata_dir: str | None,
    ) -> str:
        configured = [item.strip() for item in settings.ocr_language.split("+") if item.strip()]
        available = set(
            self._available_tesseract_languages(
                pytesseract_module=pytesseract_module,
                tessdata_dir=tessdata_dir,
            )
        )
        if not configured:
            configured = ["eng"]

        supported = [language for language in configured if language in available]
        if supported:
            return "+".join(supported)

        for fallback in ("eng", "osd"):
            if fallback in available:
                return fallback

        return settings.ocr_language

    def _available_tesseract_languages(
        self,
        pytesseract_module: object,
        tessdata_dir: str | None,
    ) -> list[str]:
        try:
            languages = pytesseract_module.get_languages(config="")
        except Exception:
            return []

        return [str(language).strip() for language in languages if str(language).strip()]

    def _resolve_tessdata_dir(self) -> str | None:
        env_raw = os.getenv("TESSDATA_PREFIX", "").strip()
        if env_raw:
            env_value = str(Path(env_raw).expanduser())
            if Path(env_value).exists():
                return env_value

        local_dir = settings.ocr_data_dir
        if local_dir.exists() and any(local_dir.glob("*.traineddata")):
            return str(local_dir)

        return None

    def _resolve_tesseract_cmd(self) -> str | None:
        if settings.tesseract_cmd:
            return settings.tesseract_cmd

        discovered = shutil.which("tesseract")
        if discovered:
            return discovered

        windows_candidates = (
            Path("C:/Program Files/Tesseract-OCR/tesseract.exe"),
            Path("C:/Program Files (x86)/Tesseract-OCR/tesseract.exe"),
        )
        for candidate in windows_candidates:
            if candidate.exists():
                return str(candidate)

        return None

    def _normalize_extracted_text(self, value: str) -> str:
        value = value.replace("\r\n", "\n").replace("\r", "\n")
        value = re.sub(r"[\x00-\x08\x0B-\x1F\x7F]+", " ", value)
        value = re.sub(r"[ \t]+", " ", value)
        value = re.sub(r"\n{3,}", "\n\n", value)
        return value

    def _extract_date_candidates(
        self,
        text: str,
        document_name: str,
    ) -> list[dict[str, str | int]]:
        normalized = self._normalize_extracted_text(text)[:20000]
        candidates: list[dict[str, str | int]] = []
        seen: set[tuple[str, str]] = set()
        month_pattern = (
            r"jan(?:uary)?|januari|feb(?:ruary)?|februari|mar(?:ch)?|mars|"
            r"apr(?:il)?|may|maj|jun(?:e)?|juni|jul(?:y)?|juli|aug(?:ust)?|"
            r"sep(?:t)?(?:ember)?|oct(?:ober)?|okt(?:ober)?|nov(?:ember)?|dec(?:ember)?"
        )
        date_value_pattern = (
            rf"(?P<date>"
            rf"\d{{4}}[/-]\d{{1,2}}[/-]\d{{1,2}}"
            rf"|\d{{1,2}}[./-]\d{{1,2}}[./-]\d{{2,4}}"
            rf"|\d{{1,2}}\.?\s+(?:{month_pattern})\.?,?\s+\d{{4}}"
            rf"|(?:{month_pattern})\.?\s+\d{{1,2}}\.?,?\s+\d{{4}}"
            rf")"
        )
        contextual_patterns = (
            ("invoice_date", rf"(?:invoice\s+date|date\s+of\s+issue|issue\s+date|fakturadato)\s*[:\-]?\s*{date_value_pattern}"),
            ("effective_date", rf"(?:effective\s+date|agreement\s+date|contract\s+date)\s*[:\-]?\s*{date_value_pattern}"),
            ("signed_date", rf"(?:signed\s+on|date\s+signed|signature\s+date)\s*[:\-]?\s*{date_value_pattern}"),
            ("policy_date", rf"(?:policy\s+date)\s*[:\-]?\s*{date_value_pattern}"),
            ("due_date", rf"(?:due\s+date)\s*[:\-]?\s*{date_value_pattern}"),
            ("document_date", rf"(?:\bdate\b)\s*[:\-]?\s*{date_value_pattern}"),
        )

        for date_kind, pattern in contextual_patterns:
            for match in re.finditer(pattern, normalized, flags=re.IGNORECASE):
                raw_value = str(match.group("date")).strip()
                parsed_date = self._parse_date_value(raw_value)
                if parsed_date is None:
                    continue

                dedupe_key = (date_kind, parsed_date.isoformat())
                if dedupe_key in seen:
                    continue

                seen.add(dedupe_key)
                candidates.append(
                    {
                        "kind": date_kind,
                        "date": parsed_date.isoformat(),
                        "raw": raw_value,
                        "position": match.start(),
                    }
                )

        for match in re.finditer(r"(\d{4})[-_](\d{2})[-_](\d{2})", document_name):
            raw_value = match.group(0)
            parsed_date = self._parse_date_value(raw_value.replace("_", "-"))
            if parsed_date is None:
                continue

            dedupe_key = ("filename_date", parsed_date.isoformat())
            if dedupe_key in seen:
                continue

            seen.add(dedupe_key)
            candidates.append(
                {
                    "kind": "filename_date",
                    "date": parsed_date.isoformat(),
                    "raw": raw_value,
                    "position": 999999,
                }
            )

        return candidates

    def _parse_date_value(self, value: str) -> date | None:
        cleaned = value.strip().replace(",", "")
        year_first_match = re.fullmatch(r"(\d{4})[/-](\d{1,2})[/-](\d{1,2})", cleaned)
        if year_first_match:
            year, month, day = (int(part) for part in year_first_match.groups())
            return self._safe_date(year, month, day)

        numeric_match = re.fullmatch(r"(\d{1,2})[./-](\d{1,2})[./-](\d{2,4})", cleaned)
        if numeric_match:
            day, month, year = (int(part) for part in numeric_match.groups())
            if year < 100:
                year = 2000 + year if year < 70 else 1900 + year
            return self._safe_date(year, month, day)

        day_name_match = re.fullmatch(
            r"(\d{1,2})\.?\s+([A-Za-z]+)\.?,?\s+(\d{4})",
            cleaned,
            flags=re.IGNORECASE,
        )
        if day_name_match:
            day = int(day_name_match.group(1))
            month_name = day_name_match.group(2).lower()
            year = int(day_name_match.group(3))
            month = self.MONTH_LOOKUP.get(month_name)
            if month:
                return self._safe_date(year, month, day)

        name_day_match = re.fullmatch(
            r"([A-Za-z]+)\.?\s+(\d{1,2})\.?,?\s+(\d{4})",
            cleaned,
            flags=re.IGNORECASE,
        )
        if name_day_match:
            month_name = name_day_match.group(1).lower()
            day = int(name_day_match.group(2))
            year = int(name_day_match.group(3))
            month = self.MONTH_LOOKUP.get(month_name)
            if month:
                return self._safe_date(year, month, day)

        return None

    def _safe_date(self, year: int, month: int, day: int) -> date | None:
        try:
            return date(year, month, day)
        except ValueError:
            return None

    def _score_date_candidate(
        self,
        candidate_kind: str,
        document_type: str,
    ) -> float:
        base_scores = {
            "invoice_date": 4.0,
            "issue_date": 3.6,
            "effective_date": 4.0,
            "signed_date": 3.6,
            "policy_date": 3.4,
            "due_date": 2.2,
            "document_date": 1.8,
            "filename_date": 1.2,
        }
        score = base_scores.get(candidate_kind, 1.0)
        if document_type == "invoice" and candidate_kind in {"invoice_date", "issue_date"}:
            score += 1.5
        if document_type == "contract" and candidate_kind in {"effective_date", "signed_date"}:
            score += 1.5
        if document_type == "insurance" and candidate_kind == "policy_date":
            score += 1.0
        if candidate_kind == "due_date" and document_type != "invoice":
            score -= 0.5
        return score

    def _cleanup_ocr_text(self, value: str) -> str:
        cleaned = self._normalize_extracted_text(value)
        cleaned = cleaned.replace("“", '"').replace("”", '"').replace("’", "'")
        cleaned = cleaned.replace("•", " ").replace("●", " ").replace("■", " ")
        cleaned = cleaned.replace("Â", " ")
        cleaned = re.sub(r"[|¦]+", " ", cleaned)
        cleaned = re.sub(r"[~^]{2,}", " ", cleaned)
        cleaned = re.sub(r"[._-]{4,}", " ", cleaned)
        cleaned = re.sub(r"\s+/\s+", " / ", cleaned)
        cleaned = re.sub(r"\bInaccordance\b", "In accordance", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bIn\s+accordance\b", "In accordance", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bFirst\s+Name[2z]\b[:.]?", "First Name:", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bLast\s+Nami\b[:.]?", "Last Name:", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bLast\s+Nam\b[:.]?", "Last Name:", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bDate\s+of\s+birth\b[:.]?", "Date of birth:", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bNo\.?\s*of\s*the\b", "No. of the", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bTel\.?\s*(?:nurnber|num ber|number)\b[:.]?", "Tel. number:", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bE-?\s*mail\s+address\b[:.]?", "Email address:", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bDoctor\s+and\s+medicine\s+expenses\b", "Doctor and medicine expenses", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bPolicy\s+holder\b", "Policy holder", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bPersonal\s+identity\s+number\b", "Personal identity number", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"([A-Za-z]{3,})\s{2,}([A-Za-z]{3,})", r"\1 \2", cleaned)

        cleaned_lines: list[str] = []
        for raw_line in cleaned.splitlines():
            line = re.sub(r"\s{2,}", " ", raw_line).strip()
            if not line:
                if cleaned_lines and cleaned_lines[-1] != "":
                    cleaned_lines.append("")
                continue

            alpha_numeric = len(re.findall(r"[A-Za-z0-9]", line))
            if alpha_numeric == 0 and len(line) <= 4:
                continue
            if alpha_numeric <= 2 and len(line) <= 6:
                continue

            noisy_characters = len(re.findall(r"[^A-Za-z0-9\s.,;:!?()/%$&@#'\"+-]", line))
            if noisy_characters >= max(4, alpha_numeric) and len(line) <= 24:
                continue

            cleaned_lines.append(line)

        return "\n".join(cleaned_lines).strip()

    def _normalize_line(self, value: str) -> str:
        return self._normalize_extracted_text(value).strip()

    def _signal_token_counts(self, text: str) -> Counter[str]:
        tokens = re.findall(r"[A-Za-z0-9][A-Za-z0-9&'/-]{2,}", text.lower())
        counts: Counter[str] = Counter()
        for token in tokens:
            normalized = token.strip(".,:;()[]{}")
            if len(normalized) < 3:
                continue
            if normalized in self.SIGNAL_STOPWORDS:
                continue
            if normalized.isdigit():
                continue
            counts[normalized] += 1
        return counts

    def _signal_phrase_counts(self, text: str) -> Counter[str]:
        words = re.findall(r"[A-Za-z0-9][A-Za-z0-9&'/-]{2,}", text)
        normalized_words = [
            word.strip(".,:;()[]{}")
            for word in words
            if word.strip(".,:;()[]{}")
        ]
        counts: Counter[str] = Counter()
        for index in range(len(normalized_words) - 1):
            left = normalized_words[index]
            right = normalized_words[index + 1]
            if (
                left.lower() in self.SIGNAL_STOPWORDS
                or right.lower() in self.SIGNAL_STOPWORDS
            ):
                continue
            phrase = f"{left} {right}"
            normalized_phrase = self._signal_key(phrase)
            if len(normalized_phrase) < 7:
                continue
            counts[phrase] += 1
        return counts

    def _signal_key(self, value: str) -> str:
        normalized = self._normalize_extracted_text(value).lower()
        normalized = normalized.replace("&", " and ")
        normalized = re.sub(r"[^a-z0-9]+", " ", normalized)
        normalized = re.sub(r"\s+", " ", normalized).strip()
        if len(normalized) < 2:
            return ""
        return normalized

    def _line_entity_candidate(
        self,
        value: str,
        document_type: str | None = None,
    ) -> str | None:
        line = self._normalize_extracted_text(value).strip()
        if not line or len(line) > 90:
            return None
        if "@" in line or "www." in line.lower():
            return None
        if re.search(r"\b(?:invoice|faktura|date|page|telefon|phone|email|e-post|bank|vat|org(?:\.|anization)?)\b", line, flags=re.IGNORECASE):
            return None

        word_count = len(line.split())
        if word_count == 0 or word_count > 7:
            return None

        if any(
            re.search(rf"\b{re.escape(suffix)}\b", line, flags=re.IGNORECASE)
            for suffix in self.COMPANY_SUFFIXES
        ):
            normalized_with_suffix = self._normalize_entity_name(line)
            if normalized_with_suffix and not self._is_noise_entity(normalized_with_suffix):
                return normalized_with_suffix

        if (
            document_type == "invoice"
            and word_count >= 2
            and re.search(r"[A-Z]{2,}|\d{2,}", line)
        ):
            normalized = self._normalize_entity_name(line)
            if (
                normalized
                and not self._is_noise_entity(normalized)
                and sum(char.isdigit() for char in normalized) <= max(4, len(normalized) // 4)
            ):
                return normalized

        return None

    def _normalize_entity_name(self, value: str) -> str | None:
        candidate = self._normalize_extracted_text(value)
        candidate = re.sub(r"^(?:vendor|supplier|seller|company|customer|bill to|ship to|invoice to|leverant[oö]r|kund)\s*[:\-]\s*", "", candidate, flags=re.IGNORECASE)
        candidate = candidate.strip(" .,:;|/-")
        candidate = re.sub(r"\s+", " ", candidate)
        if len(candidate) < 3:
            return None
        if re.fullmatch(r"[\d\W_]+", candidate):
            return None

        words: list[str] = []
        for token in candidate.split():
            clean_token = token.strip(".,:;()[]{}")
            if not clean_token:
                continue
            lowered = clean_token.lower()
            if lowered in self.COMPANY_SUFFIXES:
                words.append(self.COMPANY_SUFFIXES[lowered])
            elif clean_token.isdigit():
                words.append(clean_token)
            elif re.fullmatch(r"[A-Z]{2,}", clean_token):
                words.append(clean_token.title())
            elif re.fullmatch(r"[A-Za-z][A-Za-z0-9&'-]*", clean_token):
                words.append(clean_token[0].upper() + clean_token[1:].lower())
            else:
                words.append(clean_token)

        normalized = " ".join(words).strip()
        if len(normalized) < 3:
            return None
        if re.search(r"\b(?:invoice|faktura|receipt|roadmap|architecture|report|document|page)\b", normalized, flags=re.IGNORECASE):
            return None
        return normalized

    def _is_noise_entity(self, value: str) -> bool:
        normalized = re.sub(r"[^a-z0-9]+", " ", value.lower()).strip()
        if not normalized:
            return True
        if normalized in self.ENTITY_NOISE_TERMS:
            return True
        digit_ratio = sum(character.isdigit() for character in normalized) / max(len(normalized), 1)
        if digit_ratio > 0.22:
            return True
        if re.search(r"\b(?:usd|eur|sek|dkk|total|subtotal|vat|address|road|street|district|china|sweden)\b", normalized):
            return True
        if len(normalized.split()) == 1 and len(normalized) < 5:
            return True
        return False

    def _is_noise_signal_value(self, value: str, category: str) -> bool:
        normalized = self._signal_key(value)
        if not normalized:
            return True
        if normalized in self.ENTITY_NOISE_TERMS:
            return True
        if category in {"entity", "title", "section"} and self._is_noise_entity(value):
            return True
        if category in {"keyword", "phrase"}:
            if len(normalized) < 4:
                return True
            if sum(character.isdigit() for character in normalized) / max(len(normalized), 1) > 0.35:
                return True
        return False

    def _entity_from_domain(self, domain: str) -> str | None:
        host = domain.lower().split(".")
        if len(host) < 2:
            return None
        brand = host[-2]
        if len(brand) < 3 or brand in {"gmail", "hotmail", "outlook", "icloud"}:
            return None
        brand = brand.replace("-", " ")
        brand = re.sub(r"(\d+)", r" \1 ", brand)
        brand = re.sub(r"\s+", " ", brand).strip()
        return self._normalize_entity_name(brand)

    def _entity_match_key(self, value: str) -> str:
        normalized = value.lower()
        normalized = normalized.replace("&", " and ")
        normalized = re.sub(r"[^a-z0-9]+", " ", normalized)
        normalized = re.sub(r"\b(?:ab|aps|as|bv|gmbh|inc|llc|ltd|oy|sa|sarl)\b", " ", normalized)
        normalized = re.sub(r"\s+", " ", normalized).strip()
        return normalized

    def _is_section_heading(self, value: str) -> bool:
        heading = value.strip()
        if not heading or len(heading) > 100:
            return False

        if heading.startswith("#"):
            return True

        if re.match(r"^\d+(\.\d+)*[\).:-]?\s+[A-Za-z].{0,90}$", heading):
            return True

        if heading.endswith(":") and len(heading.split()) <= 12:
            return True

        if heading.isupper() and len(re.findall(r"[A-Z]", heading)) >= 4:
            return True

        if (
            len(heading.split()) <= 10
            and heading[-1] not in ".?!"
            and re.match(r"^[A-Z][A-Za-z0-9/&()' -]+$", heading)
        ):
            return True

        return False

    def _clean_heading(self, value: str) -> str:
        heading = value.strip().lstrip("#").strip()
        heading = re.sub(r"^\d+(\.\d+)*[\).:-]?\s*", "", heading).strip()
        return heading or "Section"

    def _is_ordered_list_item(self, value: str) -> bool:
        line = value.strip()
        return bool(
            re.match(r"^\d+[\).]\s+\S+", line)
            or re.match(r"^\d+\.\s+\S+", line)
        )
