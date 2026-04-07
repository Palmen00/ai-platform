from __future__ import annotations

import argparse
import json
import os
from dataclasses import dataclass, asdict
from datetime import datetime
from pathlib import Path
import shutil
import sys
import textwrap
import time
from typing import Any
from xml.etree.ElementTree import Element, SubElement, ElementTree

import requests
from PIL import Image, ImageDraw, ImageFont
from docx import Document
from openpyxl import Workbook
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
DEFAULT_OUTPUT_DIR = ROOT / "temp" / "upload-ocr-test"


@dataclass
class TestCase:
    key: str
    filename: str
    question: str
    expected_substring: str
    kind: str
    expects_ocr: bool = False
    content_type_hint: str | None = None


TEST_CASES: list[TestCase] = [
    TestCase(
        key="txt",
        filename="notes.txt",
        question="What is the project codename in this document?",
        expected_substring="NEBULA-FOX",
        kind="text",
    ),
    TestCase(
        key="md",
        filename="guide.md",
        question="Which deployment region is mentioned in this markdown file?",
        expected_substring="Stockholm-West",
        kind="markdown",
    ),
    TestCase(
        key="json",
        filename="profile.json",
        question="What is the support owner value in this JSON file?",
        expected_substring="Marta Linden",
        kind="json",
    ),
    TestCase(
        key="csv",
        filename="finance.csv",
        question="What is the invoice amount for row INV-77?",
        expected_substring="18450",
        kind="csv",
    ),
    TestCase(
        key="xml",
        filename="service.xml",
        question="What service port is configured in the XML file?",
        expected_substring="4317",
        kind="xml",
    ),
    TestCase(
        key="code",
        filename="worker.py",
        question="What function name returns the audit status?",
        expected_substring="build_audit_status",
        kind="code",
    ),
    TestCase(
        key="docx",
        filename="policy.docx",
        question="What policy title is written in the DOCX document?",
        expected_substring="Retention Policy Aurora",
        kind="docx",
    ),
    TestCase(
        key="xlsx",
        filename="metrics.xlsx",
        question="What is the Q2 total value in the spreadsheet?",
        expected_substring="982",
        kind="xlsx",
    ),
    TestCase(
        key="pptx",
        filename="roadmap.pptx",
        question="What milestone name appears in the presentation?",
        expected_substring="Orion Launch",
        kind="pptx",
    ),
    TestCase(
        key="png_ocr",
        filename="scan-image.png",
        question="What access code appears in the scanned image?",
        expected_substring="AURORA-17",
        kind="image",
        expects_ocr=True,
    ),
    TestCase(
        key="pdf_ocr",
        filename="scan-pdf.pdf",
        question="What incident code appears in the scanned PDF?",
        expected_substring="INC-2048",
        kind="pdf",
        expects_ocr=True,
    ),
]


class TestFailure(RuntimeError):
    pass


def _now_stamp() -> str:
    return datetime.now().strftime("%Y%m%d-%H%M%S")


def _load_font(size: int):
    candidates = [
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("C:/Windows/Fonts/calibri.ttf"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def _write_text_file(path: Path, body: str) -> None:
    path.write_text(body, encoding="utf-8")


def _build_fixtures(base_dir: Path, run_prefix: str) -> dict[str, Path]:
    if base_dir.exists():
        shutil.rmtree(base_dir)
    base_dir.mkdir(parents=True, exist_ok=True)

    files: dict[str, Path] = {}

    files["txt"] = base_dir / f"{run_prefix}-notes.txt"
    _write_text_file(
        files["txt"],
        textwrap.dedent(
            """
            Project codename: NEBULA-FOX
            Environment: validation
            Summary: This plain text file exists to verify text upload and retrieval.
            """
        ).strip()
        + "\n",
    )

    files["md"] = base_dir / f"{run_prefix}-guide.md"
    _write_text_file(
        files["md"],
        textwrap.dedent(
            """
            # Deployment guide

            Primary deployment region: Stockholm-West

            This markdown file verifies headings, paragraphs, and retrieval grounding.
            """
        ).strip()
        + "\n",
    )

    files["json"] = base_dir / f"{run_prefix}-profile.json"
    files["json"].write_text(
        json.dumps(
            {
                "service": "local-ai-os",
                "support_owner": "Marta Linden",
                "active": True,
                "note": "JSON ingestion test payload",
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    files["csv"] = base_dir / f"{run_prefix}-finance.csv"
    _write_text_file(
        files["csv"],
        textwrap.dedent(
            """
            invoice_id,amount,status
            INV-76,12000,paid
            INV-77,18450,due
            INV-78,9300,paid
            """
        ).strip()
        + "\n",
    )

    files["xml"] = base_dir / f"{run_prefix}-service.xml"
    root = Element("service")
    SubElement(root, "name").text = "telemetry"
    SubElement(root, "port").text = "4317"
    SubElement(root, "mode").text = "grpc"
    ElementTree(root).write(files["xml"], encoding="utf-8", xml_declaration=True)

    files["code"] = base_dir / f"{run_prefix}-worker.py"
    _write_text_file(
        files["code"],
        textwrap.dedent(
            """
            def build_audit_status() -> str:
                return "audit-ok"


            def helper() -> str:
                return "worker-ready"
            """
        ).strip()
        + "\n",
    )

    files["docx"] = base_dir / f"{run_prefix}-policy.docx"
    doc = Document()
    doc.add_heading("Retention Policy Aurora", level=1)
    doc.add_paragraph("Owner: Compliance Team")
    doc.add_paragraph("Purpose: Verify DOCX extraction and retrieval.")
    doc.save(files["docx"])

    files["xlsx"] = base_dir / f"{run_prefix}-metrics.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Metrics"
    ws.append(["Quarter", "Total"])
    ws.append(["Q1", 811])
    ws.append(["Q2", 982])
    ws.append(["Q3", 1044])
    wb.save(files["xlsx"])

    files["pptx"] = base_dir / f"{run_prefix}-roadmap.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Milestone: Orion Launch\nOwner: Platform Team"
    presentation.save(files["pptx"])

    font = _load_font(48)
    image = Image.new("RGB", (1400, 900), color="white")
    draw = ImageDraw.Draw(image)
    draw.multiline_text(
        (80, 120),
        "Scanned access code\nAURORA-17\nOCR validation image",
        fill="black",
        font=font,
        spacing=18,
    )
    files["png_ocr"] = base_dir / f"{run_prefix}-scan-image.png"
    image.save(files["png_ocr"])

    pdf_image = Image.new("RGB", (1500, 1000), color="white")
    pdf_draw = ImageDraw.Draw(pdf_image)
    pdf_draw.multiline_text(
        (90, 150),
        "Incident report\nCode: INC-2048\nScanned PDF OCR validation",
        fill="black",
        font=font,
        spacing=20,
    )
    files["pdf_ocr"] = base_dir / f"{run_prefix}-scan-pdf.pdf"
    pdf_image.save(files["pdf_ocr"], "PDF", resolution=150.0)

    return files


def _contains_expected(text: str, expected: str) -> bool:
    return expected.lower() in text.lower()


def _coerce_bool(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    return str(value).strip().lower() in {"1", "true", "yes", "on"}


def _fetch_json(session: requests.Session, method: str, url: str, **kwargs):
    response = session.request(method, url, timeout=kwargs.pop("timeout", 60), **kwargs)
    if not response.ok:
        raise TestFailure(f"{method} {url} failed: {response.status_code} {response.text[:500]}")
    return response.json()


def _login(session: requests.Session, base_url: str, username: str, password: str) -> None:
    payload = {"username": username, "password": password}
    response = session.post(
        f"{base_url}/auth/login",
        json=payload,
        timeout=30,
    )
    if not response.ok:
        raise TestFailure(f"Login failed: {response.status_code} {response.text[:500]}")


def _pick_model(session: requests.Session, base_url: str) -> str | None:
    response = session.get(f"{base_url}/models", timeout=30)
    if not response.ok:
        return None
    payload = response.json()
    models = payload.get("models") or []
    if not models:
        return None
    preferred = [
        model
        for model in models
        if model.get("capability") == "chat" and model.get("installed")
    ]
    if preferred:
        return preferred[0].get("name")

    preferred = [model for model in models if model.get("capability") == "chat"]
    if preferred:
        return preferred[0].get("name")

    preferred = [model for model in models if model.get("available")]
    chosen = preferred[0] if preferred else models[0]
    return chosen.get("name")


def _list_documents(session: requests.Session, base_url: str) -> list[dict[str, Any]]:
    payload = _fetch_json(
        session,
        "GET",
        f"{base_url}/documents",
        params={"limit": 500, "offset": 0},
    )
    return payload.get("documents", [])


def _upload_document(
    session: requests.Session,
    base_url: str,
    case: TestCase,
    path: Path,
) -> dict[str, Any]:
    with path.open("rb") as handle:
        files = {
            "file": (
                path.name,
                handle,
                case.content_type_hint or "application/octet-stream",
            )
        }
        response = session.post(
            f"{base_url}/documents/upload",
            files=files,
            timeout=120,
        )
    if not response.ok:
        raise TestFailure(
            f"Upload failed for {path.name}: {response.status_code} {response.text[:500]}"
        )
    payload = response.json()
    return payload["document"]


def _wait_for_document(
    session: requests.Session,
    base_url: str,
    document_id: str,
    timeout_seconds: int = 240,
) -> dict[str, Any]:
    deadline = time.time() + timeout_seconds
    last_seen: dict[str, Any] | None = None
    while time.time() < deadline:
        for document in _list_documents(session, base_url):
            if document.get("id") == document_id:
                last_seen = document
                processing_status = document.get("processing_status")
                indexing_status = document.get("indexing_status")
                if processing_status in {"processed", "failed"} and indexing_status in {
                    "indexed",
                    "failed",
                    "pending",
                }:
                    if processing_status == "processed" and indexing_status == "pending":
                        break
                    return document
        time.sleep(3)
    raise TestFailure(
        "Timed out waiting for document "
        f"{document_id}. Last state: {json.dumps(last_seen or {}, ensure_ascii=False)}"
    )


def _get_preview(session: requests.Session, base_url: str, document_id: str) -> dict[str, Any]:
    payload = _fetch_json(
        session,
        "GET",
        f"{base_url}/documents/{document_id}/preview",
        params={"chunk": 0},
        timeout=60,
    )
    return payload["preview"]


def _ask_document_question(
    session: requests.Session,
    base_url: str,
    model: str | None,
    document_id: str,
    question: str,
) -> dict[str, Any]:
    body = {
        "message": question,
        "model": model,
        "history": [],
        "conversation_id": None,
        "document_ids": [document_id],
        "persist_conversation": False,
    }
    return _fetch_json(session, "POST", f"{base_url}/chat", json=body, timeout=180)


def _delete_document(session: requests.Session, base_url: str, document_id: str) -> None:
    response = session.delete(f"{base_url}/documents/{document_id}", timeout=30)
    if response.status_code not in {200, 204}:
        raise TestFailure(
            f"Delete failed for {document_id}: {response.status_code} {response.text[:500]}"
        )


def _render_markdown_report(
    *,
    run_id: str,
    base_url: str,
    model: str | None,
    generated_at: str,
    results: list[dict[str, Any]],
) -> str:
    lines = [
        f"# Upload and OCR E2E Test Report {generated_at}",
        "",
        f"- Run ID: `{run_id}`",
        f"- API base URL: `{base_url}`",
        f"- Chat model: `{model or 'default'}`",
        f"- Total cases: `{len(results)}`",
        f"- Passed: `{sum(1 for item in results if item['overall_pass'])}`",
        f"- Failed: `{sum(1 for item in results if not item['overall_pass'])}`",
        "",
        "## Results",
        "",
        "| File | Kind | Upload | Processing | OCR | Preview | Chat | Overall |",
        "| --- | --- | --- | --- | --- | --- | --- | --- |",
    ]

    for item in results:
        lines.append(
            "| "
            + " | ".join(
                [
                    item["filename"],
                    item["kind"],
                    "ok" if item["upload_pass"] else "fail",
                    "ok" if item["processing_pass"] else "fail",
                    item["ocr_summary"],
                    "ok" if item["preview_pass"] else "fail",
                    "ok" if item["chat_pass"] else "fail",
                    "ok" if item["overall_pass"] else "fail",
                ]
            )
            + " |"
        )

    lines.extend(["", "## Detailed findings", ""])
    for item in results:
        lines.extend(
            [
                f"### {item['filename']}",
                "",
                f"- Kind: `{item['kind']}`",
                f"- Document ID: `{item.get('document_id') or 'n/a'}`",
                f"- Processing status: `{item.get('processing_status')}`",
                f"- Indexing status: `{item.get('indexing_status')}`",
                f"- OCR status: `{item.get('ocr_status')}`",
                f"- OCR used: `{item.get('ocr_used')}`",
                f"- Preview pass: `{item['preview_pass']}`",
                f"- Chat pass: `{item['chat_pass']}`",
                f"- Expected substring: `{item['expected_substring']}`",
                f"- Chat answer: `{item.get('chat_reply', '').strip()}`",
                f"- Retrieved sources: `{item.get('source_names', [])}`",
                f"- Notes: {item.get('notes') or 'None'}",
                "",
            ]
        )

    return "\n".join(lines)


def run(args: argparse.Namespace) -> int:
    run_id = _now_stamp()
    fixture_dir = args.output_dir / f"fixtures-{run_id}"
    report_json = args.output_dir / f"upload-ocr-report-{run_id}.json"
    report_md = args.output_dir / f"upload-ocr-report-{run_id}.md"
    args.output_dir.mkdir(parents=True, exist_ok=True)

    fixtures = _build_fixtures(fixture_dir, run_id)

    session = requests.Session()
    session.headers.update({"Accept": "application/json"})

    _login(session, args.base_url, args.username, args.password)
    model = _pick_model(session, args.base_url)

    results: list[dict[str, Any]] = []
    failures: list[str] = []

    for case in TEST_CASES:
        path = fixtures[case.key]
        result: dict[str, Any] = {
            "key": case.key,
            "filename": path.name,
            "kind": case.kind,
            "expected_substring": case.expected_substring,
            "upload_pass": False,
            "processing_pass": False,
            "preview_pass": False,
            "chat_pass": False,
            "overall_pass": False,
            "ocr_summary": "not_checked",
            "notes": "",
        }
        try:
            uploaded = _upload_document(session, args.base_url, case, path)
            result["upload_pass"] = True
            result["document_id"] = uploaded["id"]

            document = _wait_for_document(session, args.base_url, uploaded["id"])
            result["processing_status"] = document.get("processing_status")
            result["indexing_status"] = document.get("indexing_status")
            result["ocr_status"] = document.get("ocr_status")
            result["ocr_used"] = _coerce_bool(document.get("ocr_used"))
            result["ocr_error"] = document.get("ocr_error")
            result["processing_error"] = document.get("processing_error")
            result["indexing_error"] = document.get("indexing_error")
            result["processing_pass"] = (
                document.get("processing_status") == "processed"
                and document.get("indexing_status") == "indexed"
            )
            result["ocr_summary"] = (
                f"{document.get('ocr_status')} / used={document.get('ocr_used')}"
            )

            preview = _get_preview(session, args.base_url, uploaded["id"])
            preview_text = preview.get("extracted_text") or ""
            result["preview_excerpt"] = preview_text[:300]
            result["preview_pass"] = _contains_expected(preview_text, case.expected_substring)

            chat = _ask_document_question(
                session,
                args.base_url,
                model,
                uploaded["id"],
                case.question,
            )
            reply = chat.get("reply") or ""
            sources = chat.get("sources") or []
            result["chat_reply"] = reply
            result["source_names"] = [source.get("document_name") for source in sources]
            result["chat_pass"] = _contains_expected(reply, case.expected_substring) or any(
                _contains_expected(source.get("excerpt") or "", case.expected_substring)
                for source in sources
            )

            if case.expects_ocr:
                result["ocr_pass"] = (
                    result["ocr_used"] is True
                    and str(result.get("ocr_status")) in {"used", "fallback_used", "processed"}
                )
            else:
                result["ocr_pass"] = True

            result["overall_pass"] = all(
                [
                    result["upload_pass"],
                    result["processing_pass"],
                    result["preview_pass"],
                    result["chat_pass"],
                    result.get("ocr_pass", True),
                ]
            )
            if not result["overall_pass"]:
                result["notes"] = "One or more validation steps failed."
                failures.append(path.name)
        except Exception as exc:  # noqa: BLE001
            result["notes"] = f"{exc.__class__.__name__}: {exc}"
            failures.append(path.name)
        finally:
            results.append(result)

    invalid_case = {
        "filename": f"{run_id}-blocked.exe",
        "negative_upload_rejected": False,
        "response_detail": "",
    }
    invalid_path = fixture_dir / invalid_case["filename"]
    invalid_path.write_bytes(b"MZ-test")
    with invalid_path.open("rb") as handle:
        response = session.post(
            f"{args.base_url}/documents/upload",
            files={"file": (invalid_path.name, handle, "application/octet-stream")},
            timeout=60,
        )
    invalid_case["negative_upload_rejected"] = response.status_code == 422
    invalid_case["response_detail"] = response.text[:500]

    generated_at = datetime.now().isoformat(timespec="seconds")
    payload = {
        "generated_at": generated_at,
        "run_id": run_id,
        "base_url": args.base_url,
        "model": model,
        "results": results,
        "negative_case": invalid_case,
    }
    report_json.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")
    report_md.write_text(
        _render_markdown_report(
            run_id=run_id,
            base_url=args.base_url,
            model=model,
            generated_at=generated_at,
            results=results,
        )
        + "\n\n## Negative case\n\n"
        + f"- Unsupported `.exe` rejected: `{invalid_case['negative_upload_rejected']}`\n"
        + f"- Response: `{invalid_case['response_detail']}`\n",
        encoding="utf-8",
    )

    print(f"JSON report: {report_json}")
    print(f"Markdown report: {report_md}")
    print(
        "Summary: "
        f"{sum(1 for item in results if item['overall_pass'])}/{len(results)} passed; "
        f"unsupported upload rejected={invalid_case['negative_upload_rejected']}"
    )

    if args.cleanup:
        for item in results:
            document_id = item.get("document_id")
            if document_id:
                try:
                    _delete_document(session, args.base_url, document_id)
                except Exception as exc:  # noqa: BLE001
                    print(f"Cleanup warning for {document_id}: {exc}", file=sys.stderr)

    return 1 if failures or not invalid_case["negative_upload_rejected"] else 0


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Run upload/OCR end-to-end tests.")
    parser.add_argument(
        "--base-url",
        default=os.getenv("LOCAL_AI_OS_BASE_URL", "http://192.168.1.105:8000"),
    )
    parser.add_argument(
        "--username",
        default=os.getenv("LOCAL_AI_OS_USERNAME", "Admin"),
    )
    parser.add_argument(
        "--password",
        default=os.getenv("LOCAL_AI_OS_PASSWORD", "password"),
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=DEFAULT_OUTPUT_DIR,
    )
    parser.add_argument(
        "--cleanup",
        action="store_true",
        help="Delete uploaded test documents after the run.",
    )
    return parser.parse_args()


if __name__ == "__main__":
    raise SystemExit(run(parse_args()))
