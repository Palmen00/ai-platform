from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation


REPO_ROOT = Path(__file__).resolve().parents[2]
FIXTURE_DIR = REPO_ROOT / "backend" / "evals" / "fixtures" / "synthetic"


def build_docx_fixture() -> None:
    path = FIXTURE_DIR / "BlueHarbor_SharePoint_Knowledge_Runbook.docx"
    document = Document()
    document.add_heading("BlueHarbor SharePoint Knowledge Runbook", level=1)
    document.add_paragraph(
        "This runbook describes how BlueHarbor syncs SharePoint knowledge libraries into the local AI knowledge index."
    )
    document.add_heading("OAuth token refresh", level=2)
    document.add_paragraph(
        "The sync worker refreshes the OAuth access token before calling the Microsoft Graph API for large SharePoint sync jobs."
    )
    document.add_heading("Incremental sync", level=2)
    document.add_paragraph(
        "Incremental sync only uploads changed files and preserves the existing chunk ids for unchanged documents."
    )
    document.add_heading("Supported file types", level=2)
    document.add_paragraph(
        "The current SharePoint lane targets PDF, DOCX, XLSX, PPTX, code files, and YAML configuration."
    )
    document.save(path)


def build_xlsx_fixture() -> None:
    path = FIXTURE_DIR / "Northwind_M365_License_Register.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Licenses"
    sheet.append(["Department", "Product", "Seats", "Notes"])
    sheet.append(["Operations", "Microsoft 365 Business Premium", 42, "Includes SharePoint connector access"])
    sheet.append(["Security", "Power BI Pro", 12, "Used for reporting dashboards"])
    sheet.append(["Support", "Exchange Online", 18, "Mailbox-only users"])

    sheet2 = workbook.create_sheet("Renewals")
    sheet2.append(["Item", "Renewal Date", "Owner"])
    sheet2.append(["Microsoft 365 Business Premium", "2026-11-30", "Northwind IT"])
    sheet2.append(["Power BI Pro", "2026-09-15", "Northwind IT"])
    workbook.save(path)


def build_pptx_fixture() -> None:
    path = FIXTURE_DIR / "Solstice_SharePoint_Rollout_Briefing.pptx"
    presentation = Presentation()

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Solstice SharePoint Rollout Briefing"
    slide.placeholders[1].text = (
        "Phased migration plan for the Solstice Logistics knowledge rollout."
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Rollout phases"
    slide.placeholders[1].text = (
        "Phase 1 covers contracts and policies. Phase 2 covers invoices, code repositories, and project documentation."
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Search goals"
    slide.placeholders[1].text = (
        "Users should be able to search SharePoint files by document type, date, company, and project name."
    )

    presentation.save(path)


def main() -> int:
    FIXTURE_DIR.mkdir(parents=True, exist_ok=True)
    build_docx_fixture()
    build_xlsx_fixture()
    build_pptx_fixture()
    print("Generated Office fixtures in", FIXTURE_DIR)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
